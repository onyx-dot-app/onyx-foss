"""Unit tests for the pptx skill's layout linter (scripts/lint.py).

Fixture decks are generated in-test with python-pptx — one per defect class
plus a clean deck — and thresholds are validated in both directions: each
defect is caught, and the clean deck produces zero findings.
"""

import importlib.util
import sys
from pathlib import Path
from types import ModuleType

import pytest

pptx = pytest.importorskip("pptx")
pytest.importorskip("PIL")

from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import MSO_AUTO_SIZE  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.util import Pt  # noqa: E402

LINT_PATH = (
    Path(__file__).parents[4]
    / "onyx"
    / "skills"
    / "builtin"
    / "pptx"
    / "scripts"
    / "lint.py"
)


def _load_lint() -> ModuleType:
    spec = importlib.util.spec_from_file_location("pptx_lint", LINT_PATH)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules["pptx_lint"] = module
    spec.loader.exec_module(module)
    return module


lint = _load_lint()


def _new_deck() -> "pptx.presentation.Presentation":
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _add_blank_slide(prs: "pptx.presentation.Presentation"):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size_pt: int = 18,
    font_name: str = "Inter",
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    # python-pptx textboxes default to no-wrap + shape-grows-to-fit; model the
    # common fixed-box wrapped layout instead.
    box.text_frame.word_wrap = True
    box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    para = box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.name = font_name
    return box


def _findings_by_check(prs, check: str) -> list:
    findings, _ = lint.lint_presentation(prs)
    return [f for f in findings if f.check == check]


def test_clean_deck_has_no_findings() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 0.5, 0.5, 9.0, 1.0, "Quarterly Review", size_pt=40)
    _add_textbox(
        slide,
        0.5,
        2.0,
        8.0,
        3.0,
        "Revenue grew 12% quarter over quarter.",
        size_pt=16,
    )
    findings, _ = lint.lint_presentation(prs)
    assert findings == []


def test_off_slide_shape_is_error() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    # 2 of 3 inches hang off the left edge.
    shape = _add_textbox(slide, -2.0, 2.0, 3.0, 1.0, "Off slide text")
    findings = _findings_by_check(prs, "BOUNDS")
    assert len(findings) == 1
    finding = findings[0]
    assert finding.severity == "ERROR"
    assert finding.slide == 1
    assert finding.shape_id == shape.shape_id


def test_sub_margin_text_is_warned() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    shape = _add_textbox(slide, 0.1, 2.0, 4.0, 1.0, "Too close to the edge")
    findings = _findings_by_check(prs, "MARGIN")
    assert len(findings) == 1
    finding = findings[0]
    assert finding.severity == "WARN"
    assert finding.slide == 1
    assert finding.shape_id == shape.shape_id
    assert "left" in finding.detail


def test_dense_profile_relaxes_margin_warning() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 0.3, 2.0, 4.0, 1.0, "Consulting-dense label")
    standard = [f for f in lint.lint_presentation(prs)[0] if f.check == "MARGIN"]
    dense = [
        f
        for f in lint.lint_presentation(prs, profile="dense")[0]
        if f.check == "MARGIN"
    ]
    assert len(standard) == 1 and standard[0].severity == "WARN"
    assert dense == []


def test_margin_finding_reports_profile_nominal_limit() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    # 0.1" from the left: inside both the standard (0.5") and dense (0.25") margins.
    _add_textbox(slide, 0.1, 2.0, 4.0, 1.0, "Edge-hugging label")
    standard = [f for f in lint.lint_presentation(prs)[0] if f.check == "MARGIN"]
    dense = [
        f
        for f in lint.lint_presentation(prs, profile="dense")[0]
        if f.check == "MARGIN"
    ]
    assert len(standard) == 1 and '0.5"' in standard[0].detail
    assert len(dense) == 1 and '0.25"' in dense[0].detail


def test_unknown_profile_raises_value_error() -> None:
    prs = _new_deck()
    with pytest.raises(ValueError, match="unknown profile"):
        lint.lint_presentation(prs, profile="premium")


def test_full_bleed_shape_exempt_from_margins_and_bounds() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    # Full-slide background with slight bleed — a standard design pattern.
    shape = slide.shapes.add_shape(
        1, Inches(-0.05), Inches(-0.05), Inches(13.45), Inches(7.6)
    )
    shape.text_frame.text = "SECTION"
    assert _findings_by_check(prs, "BOUNDS") == []
    assert _findings_by_check(prs, "MARGIN") == []


def test_overflowing_textbox_is_caught() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    long_text = (
        "This is a very long paragraph of body copy that cannot possibly fit "
        "inside a half-inch-tall text box because it will wrap onto many lines "
        "once the eighteen point font is measured against the narrow box width "
        "and the estimated height greatly exceeds the available space."
    )
    shape = _add_textbox(slide, 2.0, 2.0, 2.0, 0.5, long_text, size_pt=18)
    findings = _findings_by_check(prs, "OVERFLOW")
    assert len(findings) == 1
    finding = findings[0]
    assert finding.severity == "ERROR"
    assert finding.shape_id == shape.shape_id


def test_short_text_in_roomy_box_does_not_overflow() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 2.0, 2.0, 8.0, 2.0, "Short line.", size_pt=16)
    assert _findings_by_check(prs, "OVERFLOW") == []


def test_overlapping_text_frames_is_error() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    a = _add_textbox(slide, 2.0, 2.0, 4.0, 2.0, "First text block")
    _add_textbox(slide, 4.0, 3.0, 4.0, 2.0, "Second text block")
    findings = _findings_by_check(prs, "OVERLAP")
    assert len(findings) == 1
    finding = findings[0]
    assert finding.severity == "ERROR"
    assert finding.shape_id == a.shape_id
    assert "intersect" in finding.detail


def test_contained_text_frames_are_allowed() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 2.0, 2.0, 6.0, 3.0, "Outer text")
    _add_textbox(slide, 3.0, 3.0, 2.0, 1.0, "Inner text")
    assert _findings_by_check(prs, "OVERLAP") == []


def test_text_spilling_past_container_card_is_warned() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    # Card (no text) with a text box centered on it but spilling out the right.
    slide.shapes.add_shape(1, Inches(2.0), Inches(2.0), Inches(3.0), Inches(2.0))
    text = _add_textbox(slide, 2.2, 2.5, 4.0, 1.0, "Spills past the card edge")
    findings = _findings_by_check(prs, "OVERLAP")
    assert len(findings) == 1
    finding = findings[0]
    assert finding.severity == "WARN"
    assert finding.shape_id == text.shape_id
    assert "container" in finding.detail


def test_low_contrast_text_is_error() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    shape = slide.shapes.add_shape(
        1, Inches(2.0), Inches(2.0), Inches(5.0), Inches(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    shape.line.fill.background()
    para = shape.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "Barely visible"
    run.font.size = Pt(18)
    run.font.name = "Inter"
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    findings = _findings_by_check(prs, "CONTRAST")
    assert len(findings) == 1
    finding = findings[0]
    assert finding.severity == "ERROR"
    assert finding.shape_id == shape.shape_id


def test_high_contrast_text_passes() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    shape = _add_textbox(slide, 2.0, 2.0, 5.0, 1.0, "Crisp dark text")
    run = shape.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor(0x21, 0x21, 0x21)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
    assert _findings_by_check(prs, "CONTRAST") == []


def test_unknown_font_is_warned_once() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    _add_textbox(slide, 2.0, 2.0, 5.0, 1.0, "Fancy", font_name="Papyrus Deluxe")
    _add_textbox(slide, 2.0, 4.0, 5.0, 1.0, "Fancy again", font_name="Papyrus Deluxe")
    findings = _findings_by_check(prs, "FONT")
    assert len(findings) == 1
    finding = findings[0]
    assert finding.severity == "WARN"
    assert "Papyrus Deluxe" in finding.detail


def test_errors_drive_summary_counts() -> None:
    prs = _new_deck()
    slide = _add_blank_slide(prs)
    _add_textbox(slide, -2.0, 2.0, 3.0, 1.0, "Off slide")
    findings, _ = lint.lint_presentation(prs)
    assert any(f.severity == "ERROR" for f in findings)
    rendered = findings[0].render()
    assert rendered.startswith("slide 1 | shape ")
    assert "| ERROR " in rendered
